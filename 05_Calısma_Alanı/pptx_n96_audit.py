# -*- coding: utf-8 -*-
"""
UYIK_2026_BES.pptx — N=96 / B2 / B3 ÇAPRAZ DENETİM
Sunumdaki tüm sayısal değerleri CSV ground-truth ile karşılaştır
"""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.util import Inches
from pathlib import Path

BASE = Path(r"c:\Users\Kurt\Desktop\Proje")
pptx_path = BASE / "01_Savunma_ve_Ana_Metinler" / "UYIK_2026_BES.pptx"

prs = Presentation(str(pptx_path))

print("=" * 80)
print("UYIK_2026_BES.pptx — TAM İÇERİK TARAMASI")
print("=" * 80)
print(f"Toplam slayt: {len(prs.slides)}")

# Kritik anahtar kelimeleri tara
keywords = ['80.21', '84.38', '77/96', '81/96', 'N=96', 'N = 96', 'pooled', 'Pooled',
            '75.56', '84.86', '0.84', '0.857', '63', '18', 'TP', 'TN', 'FP', 'FN',
            'confusion', 'matris', 'matrix', 'seed', 'Seed',
            '78.79', 'Naive', 'naive',
            # B2 ilgili değerler
            '77.38', '69.61', '75.00', '67.86', '74.36', '60.42', '72.92',
            # Feature set
            'full', 'technical', 'closing']

for slide_idx, slide in enumerate(prs.slides):
    slide_num = slide_idx + 1
    all_text = []
    
    # Tüm shape'lerden metin çek
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    all_text.append(text)
        
        # Tablolar
        if shape.has_table:
            table = shape.table
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    cell_text = cell.text.strip()
                    if cell_text:
                        row_text.append(cell_text)
                if row_text:
                    all_text.append(" | ".join(row_text))
    
    # Kritik slaytları detaylı göster
    full_text = "\n".join(all_text)
    has_critical = any(kw in full_text for kw in ['80.21', '84.38', 'N=96', 'pooled', 'Pooled', 
                                                    'TP', 'confusion', 'matris', 'matrix',
                                                    '75.56', 'Naive', '78.79', 'seed', 'Seed',
                                                    '77.38', '69.61', '67.86'])
    
    if has_critical:
        print(f"\n{'='*80}")
        print(f"SLAYT {slide_num} — KRİTİK İÇERİK")
        print(f"{'='*80}")
        for line in all_text:
            print(f"  {line[:200]}")
    else:
        # Sadece başlık göster
        title = all_text[0] if all_text else "(boş)"
        print(f"\nSlayt {slide_num}: {title[:80]}")
