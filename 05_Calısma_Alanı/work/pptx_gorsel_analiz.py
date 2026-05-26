# -*- coding: utf-8 -*-
"""PPTX Görsel karşılaştırma + Referans analizi"""
import sys, io, os, hashlib
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
from pptx import Presentation

BASE = r'c:\Users\Kurt\Desktop\Proje'

def extract_images_info(pptx_path):
    """Extract image details from PPTX"""
    prs = Presentation(pptx_path)
    images = []
    for i, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.shape_type == 13:  # Picture
                img = shape.image
                h = hashlib.md5(img.blob).hexdigest()[:8]
                images.append({
                    'slide': i,
                    'size': len(img.blob),
                    'type': img.content_type,
                    'hash': h,
                    'width': shape.width,
                    'height': shape.height
                })
    return images

def extract_slide_content(pptx_path):
    """Full slide content extraction"""
    prs = Presentation(pptx_path)
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        imgs = 0
        tbls = 0
        notes = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    t = p.text.strip()
                    if t:
                        texts.append(t)
            if shape.shape_type == 13: imgs += 1
            if shape.has_table: tbls += 1
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        slides.append({'num': i, 'texts': texts, 'imgs': imgs, 'tbls': tbls, 'notes': notes})
    return slides

# === 1. UYIK_1 vs UYIK_2026_FINAL görsel karşılaştırması ===
print("="*70)
print("1. GORSEL KARSILASTIRMA: UYIK_1 vs UYIK_2026_FINAL")
print("="*70)

imgs1 = extract_images_info(os.path.join(BASE, 'BES_Pension_UYIK_1.pptx'))
imgs2 = extract_images_info(os.path.join(BASE, 'BES_Pension_UYIK_2026_FINAL.pptx'))

print(f"\nUYIK_1: {len(imgs1)} gorsel")
for img in imgs1:
    print(f"  S{img['slide']:02d} | {img['size']:>8,} bytes | hash={img['hash']} | {img['type']}")

print(f"\nUYIK_2026_FINAL: {len(imgs2)} gorsel")
for img in imgs2:
    print(f"  S{img['slide']:02d} | {img['size']:>8,} bytes | hash={img['hash']} | {img['type']}")

# Ortak görseller (hash ile)
hashes1 = {img['hash'] for img in imgs1}
hashes2 = {img['hash'] for img in imgs2}
common = hashes1 & hashes2
only1 = hashes1 - hashes2
only2 = hashes2 - hashes1

print(f"\nOrtak gorsel: {len(common)}")
print(f"Sadece UYIK_1: {len(only1)} gorsel")
print(f"Sadece UYIK_2026_FINAL: {len(only2)} gorsel")

# === 2. Referans PPTX (01 icindeki) yorumları ===
print("\n" + "="*70)
print("2. REFERANS PPTX YORUMLARI (01/BES_Pension_Funds_Presentation_EN_OKE_20.04.26_UYIK.pptx)")
print("="*70)

ref_path = os.path.join(BASE, '01_Savunma_ve_Ana_Metinler', 'BES_Pension_Funds_Presentation_EN_OKE_20.04.26_UYIK.pptx')
ref_slides = extract_slide_content(ref_path)
ref_imgs = extract_images_info(ref_path)

print(f"Slayt: {len(ref_slides)}, Gorsel: {len(ref_imgs)}")
for s in ref_slides:
    title = ""
    for t in s['texts']:
        if len(t) > 3 and t not in ['UYIK']:
            title = t[:80]
            break
    extra = ""
    if s['imgs']: extra += f" [{s['imgs']} gorsel]"
    if s['tbls']: extra += f" [{s['tbls']} tablo]"
    notes_preview = ""
    if s['notes']:
        notes_preview = f"\n    YORUM: {s['notes'][:200]}"
    print(f"  S{s['num']:02d}: {title}{extra}{notes_preview}")

# === 3. Ornek sunum (05) ===
print("\n" + "="*70)
print("3. ORNEK SUNUM (05/uyik-2025)")
print("="*70)

ornek_path = os.path.join(BASE, '05_Hoca_Onerlileri_Ve_Ornekler', 'uyik-2025-sunum-G.Ö-Ö.K.E-B.B.K.12052025_2000 (2).pptx')
ornek_slides = extract_slide_content(ornek_path)
ornek_imgs = extract_images_info(ornek_path)

print(f"Slayt: {len(ornek_slides)}, Gorsel: {len(ornek_imgs)}")
for s in ornek_slides:
    title = ""
    for t in s['texts']:
        if len(t) > 3:
            title = t[:80]
            break
    extra = ""
    if s['imgs']: extra += f" [{s['imgs']} gorsel]"
    if s['tbls']: extra += f" [{s['tbls']} tablo]"
    print(f"  S{s['num']:02d}: {title}{extra}")

# === 4. 04 Gorsel Portfolyo - BES uyumlu olanlar ===
print("\n" + "="*70)
print("4. 04_Gorsel_Portfolyo - THYAO'dan arindirilmis gorseller")
print("="*70)

grafik_dir = os.path.join(BASE, '04_Gorsel_Portfolyo')
thyao_keywords = ['thyao','bist','hisse','stock','closing']

for subdir in ['', 'EN_Graphics', 'TR_Grafikler']:
    d = os.path.join(grafik_dir, subdir) if subdir else grafik_dir
    print(f"\n  [{subdir or 'kok'}]:")
    for f in sorted(os.listdir(d)):
        if f.endswith('.png'):
            full = os.path.join(d, f)
            size = os.path.getsize(full)
            fname_lower = f.lower()
            # Check if THYAO-specific
            is_thyao = any(kw in fname_lower for kw in thyao_keywords)
            # Check content from filename
            bes_ok = not is_thyao  # If not THYAO-specific, it's BES-compatible
            status = "BES OK" if bes_ok else "THYAO-ONLY"
            print(f"    {f:45s} {size:>8,} bytes  [{status}]")

print("\n=== TAMAMLANDI ===")
