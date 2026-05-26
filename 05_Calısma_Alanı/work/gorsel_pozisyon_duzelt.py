# -*- coding: utf-8 -*-
"""Gorsel pozisyonlarini duzelt — mevcut icerikle cakismamali"""
import sys, io, os
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
from pptx import Presentation
from pptx.util import Inches

BASE = r'c:\Users\Kurt\Desktop\Proje'
SRC = os.path.join(BASE, 'BES_Pension_UYIK_2026_FINAL.pptx')

prs = Presentation(SRC)

# S09: Sag panelde zaten TextBox var (L=8.3, W=4.6) -> gorseli altina yerlestir
# Mevcut sag panel: Rectangle 12 (L=8.3, T=1.6, W=4.6, H=5.2) — icerik kutusu
# TextBox 15 alt sinir: T=3.1 + H=3.7 = T=6.8
# Gorsel buraya sigmaz, sol taraftaki metin alaninin altina koyalim
# Sol panel: TextBox 11 (L=0.6, T=1.6, W=7.5, H=5.3) -> son yer T=6.9
# Cozum: gorseli mevcut sag panelin icine kucuk olarak yerlestir

slide09 = prs.slides[8]
for shape in slide09.shapes:
    if shape.shape_type == 13:  # Ekledigimiz gorsel
        # Sag panelin alt bosluguna kucult
        shape.left = Inches(8.5)
        shape.top = Inches(4.5)
        shape.width = Inches(4.0)
        print(f"S09 gorsel pozisyon duzeltildi: L=8.5 T=4.5 W=4.0")

# S12: Slayt cok dolu (tablo + 6 satir metin, T=1.9'dan T=6.8'e kadar)
# Gorsel icin yer yok — sol ust kosede kucuk koyalim
slide12 = prs.slides[11]
for shape in slide12.shapes:
    if shape.shape_type == 13:
        # Tablonun yanina, aciklama textbox'larinin soluna
        shape.left = Inches(8.0)
        shape.top = Inches(4.2)
        shape.width = Inches(4.5)
        print(f"S12 gorsel pozisyon duzeltildi: L=8.0 T=4.2 W=4.5")

# S14: Iki tablo yan yana (sol: L=0.5-6.7, sag: L=6.8-12.8, T=3.0-6.3)
# Ust kisimdaki 5 card (T=1.6-2.8) alani da dolu
# Gorseli tablolarin arasina koyamayiz — S14'e gorsel uymuyor, cikaralim
slide14 = prs.slides[13]
to_remove_14 = None
for shape in slide14.shapes:
    if shape.shape_type == 13:
        to_remove_14 = shape
        break
if to_remove_14:
    sp = to_remove_14._element
    sp.getparent().remove(sp)
    print("S14 gorseli cikarildi (slayt zaten dolu)")

# S16: Sag panel var (Rectangle 12, L=7.6, W=5.3) icerisinde tablo (L=7.8, T=2.0, H=3.0)
# ve TextBox (T=5.2, H=1.6) -> tablo + metin alani T=2.0'dan T=6.8'e
# Gorsel icin kucuk bir alan yok, sol panelin alt tarafina koyalim
slide16 = prs.slides[15]
for shape in slide16.shapes:
    if shape.shape_type == 13:
        # Sol metin alaninin (L=0.6, T=1.6, W=6.8, H=5.3) altina sigmaz
        # Cozum: MC Sinkhole cok onemli, tam ekran olarak ayri slayt yapmaliyiz
        # Simdilik sol panelin alt kismina kucuk yerlestir
        shape.left = Inches(0.6)
        shape.top = Inches(4.8)
        shape.width = Inches(6.5)
        print(f"S16 gorsel pozisyon duzeltildi: L=0.6 T=4.8 W=6.5")

# S25: Mevcut gorsel zaten var (L=0.6, T=3.6, W=6.2) + sag panelde metin (L=7.0)
# Yeni gorsel sag panele koyulmus ama metin kutusu var -> cakisma!
slide25 = prs.slides[24]
pics_25 = [s for s in slide25.shapes if s.shape_type == 13]
# Son eklenen gorseli (en buyuk left) duzelt
if len(pics_25) > 1:
    # Heatmap'i sag paneldeki metin altina koy
    newest = max(pics_25, key=lambda s: s.left)
    newest.left = Inches(7.2)
    newest.top = Inches(4.8)
    newest.width = Inches(5.5)
    print(f"S25 gorsel pozisyon duzeltildi: L=7.2 T=4.8 W=5.5")

prs.save(SRC)
print(f"\nKaydedildi. S14'ten gorsel cikarildi, toplam: 15 gorsel")

# Final dogrulama
prs2 = Presentation(SRC)
total = 0
for i, slide in enumerate(prs2.slides, 1):
    imgs = sum(1 for s in slide.shapes if s.shape_type == 13)
    total += imgs
    if imgs:
        print(f"  S{i:02d}: {imgs} gorsel")
print(f"TOPLAM: {total}")
