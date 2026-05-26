# -*- coding: utf-8 -*-
"""Sorunlu slaytlarin detayli layout analizi"""
import sys, io, os
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
from pptx import Presentation
from pptx.util import Inches, Emu
from PIL import Image as PILImage

BASE = r'c:\Users\Kurt\Desktop\Proje'
prs = Presentation(os.path.join(BASE, 'BES_Pension_UYIK_2026_FINAL.pptx'))

slide_w = prs.slide_width / 914400   # 13.33"
slide_h = prs.slide_height / 914400  # 7.50"

for slide_num in [1, 2, 16, 25, 26]:
    idx = slide_num - 1
    slide = prs.slides[idx]
    print('=' * 70)
    print('SLAYT %d:' % slide_num)
    print('=' * 70)
    
    for shape in slide.shapes:
        L = shape.left / 914400 if shape.left else 0
        T = shape.top / 914400 if shape.top else 0
        W = shape.width / 914400 if shape.width else 0
        H = shape.height / 914400 if shape.height else 0
        
        stype = 'UNKNOWN'
        extra = ''
        if shape.shape_type == 13:
            stype = 'PICTURE'
            try:
                img = shape.image
                extra = ' [%s, %d bytes]' % (img.content_type, len(img.blob))
            except:
                pass
        elif shape.has_table:
            stype = 'TABLE'
            tbl = shape.table
            extra = ' [%dx%d]' % (len(tbl.rows), len(tbl.columns))
        elif shape.has_text_frame:
            stype = 'TEXT'
            txt = ''
            for p in shape.text_frame.paragraphs:
                t = p.text.strip()
                if t:
                    txt = t[:50]
                    break
            extra = ' "%s"' % txt if txt else ' (bos)'
        else:
            stype = 'SHAPE'
        
        # Tasmayı kontrol et
        overflow = ''
        if L + W > slide_w + 0.1:
            overflow = ' *** SAG TASIYOR! ***'
        if T + H > slide_h + 0.1:
            overflow = ' *** ALT TASIYOR! ***'
        if L < 0 or T < 0:
            overflow = ' *** NEGATIF POZISYON! ***'
            
        print('  %-8s L=%.2f T=%.2f W=%.2f H=%.2f%s%s' % (stype, L, T, W, H, extra, overflow))
    print()
