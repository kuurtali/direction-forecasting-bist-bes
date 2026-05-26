# -*- coding: utf-8 -*-
"""S26 footer tasmasi duzelt — tum slaytlarda kontrol et"""
import sys, io, os
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
from pptx import Presentation
from pptx.util import Inches

prs = Presentation(r'c:\Users\Kurt\Desktop\Proje\BES_Pension_UYIK_2026_FINAL.pptx')

# Footer T=7.18-7.19 olan textbox'lar tum slaytlarda var (template'in parcasi)
# Bunlar asil slide master'dan geliyor, alt sinir: T=7.19+H=0.30=7.49 < 7.50
# PowerPoint'te 7.49 < 7.50 oldugu icin tasmaz aslinda
# Ama S01'deki ozel footer kutusunu (Rectangle H=0.75) zaten kucultttuk

# S26'da gorseli biraz daha yukari kaydirip caption'a yer acalim
slide26 = prs.slides[25]
for shape in slide26.shapes:
    if shape.shape_type == 13:
        L = shape.left / 914400
        T = shape.top / 914400
        W = shape.width / 914400
        H = shape.height / 914400
        print('S26 gorsel oncesi: L=%.2f T=%.2f W=%.2f H=%.2f alt=%.2f' % (L, T, W, H, T+H))
        # Gorsel zaten OK: alt=6.00, caption T=6.43 — yer var
        
# Tum slaytlarda footer kontrolu
footer_issues = 0
for i, slide in enumerate(prs.slides, 1):
    for shape in slide.shapes:
        if shape.has_text_frame:
            T = shape.top / 914400 if shape.top else 0
            H = shape.height / 914400 if shape.height else 0
            bottom = T + H
            if bottom > 7.50:
                footer_issues += 1
                print('TASMA S%02d: T=%.2f H=%.2f alt=%.2f' % (i, T, H, bottom))

if footer_issues == 0:
    print('Hicbir slaytda tasma YOK — tum icerikler slayt sinirlarinda')
else:
    print('Toplam %d tasma tespit edildi' % footer_issues)
