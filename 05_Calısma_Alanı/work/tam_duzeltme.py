# -*- coding: utf-8 -*-
"""Tum slayt sorunlarini duzelt"""
import sys, io, os
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
from pptx import Presentation
from pptx.util import Inches, Emu

BASE = r'c:\Users\Kurt\Desktop\Proje'

# BACKUP'tan tekrar baslayalim (temiz kopya)
import shutil
BACKUP = os.path.join(BASE, 'BES_Pension_UYIK_2026_FINAL_BACKUP.pptx')
DST = os.path.join(BASE, 'BES_Pension_UYIK_2026_FINAL.pptx')
shutil.copy2(BACKUP, DST)
print("BACKUP'tan temiz kopya alindi")

prs = Presentation(DST)
# slide_w = 13.33", slide_h = 7.50"
# Header: T=0.00, H=0.35
# Footer: T=7.18, H=0.32
# Icerik alani: T=1.40 -> T=7.10 (yukseklik ~5.70")

GRAFIK = os.path.join(BASE, '04_Gorsel_Portfolyo')

# Once tum slayt layoutlarini analiz edelim
def get_layout(slide):
    """Slayttaki icerik panellerini bul"""
    panels = []
    for shape in slide.shapes:
        L = shape.left / 914400 if shape.left else 0
        T = shape.top / 914400 if shape.top else 0
        W = shape.width / 914400 if shape.width else 0
        H = shape.height / 914400 if shape.height else 0
        panels.append({'L': L, 'T': T, 'W': W, 'H': H, 'shape': shape})
    return panels

# ====================================================
# S01 TASMA DUZELTME
# ====================================================
slide01 = prs.slides[0]
print('\n--- S01 analiz ---')
for shape in slide01.shapes:
    L = shape.left / 914400 if shape.left else 0
    T = shape.top / 914400 if shape.top else 0
    W = shape.width / 914400 if shape.width else 0
    H = shape.height / 914400 if shape.height else 0
    bottom = T + H
    right = L + W
    txt = ''
    if shape.has_text_frame:
        for p in shape.text_frame.paragraphs:
            t = p.text.strip()
            if t:
                txt = t[:40]
                break
    overflow = ''
    if bottom > 7.50: overflow = ' *** ALT TASMA ***'
    if right > 13.35: overflow = ' *** SAG TASMA ***'
    print('  L=%.2f T=%.2f W=%.2f H=%.2f b=%.2f r=%.2f "%s"%s' % (L, T, W, H, bottom, right, txt, overflow))

# S01 duzelt: alt footer alani (T=6.75 H=0.75) icindeki textler
# "April 2026" T=6.80 ve "Presented at UYIK" T=7.12 -> alt = 7.47
# Bunlar footer box icinde, tasma yok teknik olarak
# AMA: Yazar kutulari: T=4.60 H=1.80 -> alt = 6.40 — TAMAM
# Baslik: T=2.20 H=1.60 -> alt = 3.80 — TAMAM
# E-posta: T=6.03 H=0.30 -> alt = 6.33 — TAMAM

# Gercek sorun: belki tum shape'lerin toplam alani 7.50"yi asiyor
# Footer'daki bos kutu: L=0.00 T=6.75 W=13.33 H=0.75 -> alt = 7.50 EXACT
# Bu tamam ama PowerPoint'te render sirasinda kucuk tasma gorunebilir
# Cozum: Footer kutusunu biraz kucult
for shape in slide01.shapes:
    T = shape.top / 914400 if shape.top else 0
    H = shape.height / 914400 if shape.height else 0
    if abs(T - 6.75) < 0.05 and H > 0.5:
        # Footer box — kucult
        shape.height = Inches(0.65)
        print('  S01: Footer kutusu kucultuldu (H=0.75 -> 0.65)')

# ====================================================
# S26 TASMA DUZELTME
# ====================================================
slide26 = prs.slides[25]
print('\n--- S26 analiz ---')
for shape in slide26.shapes:
    L = shape.left / 914400 if shape.left else 0
    T = shape.top / 914400 if shape.top else 0
    W = shape.width / 914400 if shape.width else 0
    H = shape.height / 914400 if shape.height else 0
    bottom = T + H
    txt = ''
    if shape.has_text_frame:
        for p in shape.text_frame.paragraphs:
            t = p.text.strip()
            if t:
                txt = t[:40]
                break
    overflow = ''
    if bottom > 7.20: overflow = ' *** TASMA ***'
    if shape.shape_type == 13:
        print('  GORSEL: L=%.2f T=%.2f W=%.2f H=%.2f b=%.2f%s' % (L, T, W, H, bottom, overflow))
    elif shape.has_table:
        print('  TABLO: L=%.2f T=%.2f W=%.2f H=%.2f b=%.2f%s' % (L, T, W, H, bottom, overflow))
    elif txt:
        print('  TEXT: L=%.2f T=%.2f W=%.2f H=%.2f b=%.2f "%s"%s' % (L, T, W, H, bottom, txt, overflow))

# S26: Original gorsel (L=0.55 T=1.60 W=6.20 H=4.80) alt=6.40
# Caption: T=6.43, H=0.30, alt=6.73 — TAMAM
# Bottom text bar: T=6.45 H=0.55, alt=7.00 — TAMAM
# Alt bar text: T=6.48 H=0.48, alt=6.96 — TAMAM
# Interpretation text: T=4.05 H=2.30, alt=6.35 — TAMAM
# Sag alt panel: footer T=7.18 H=0.32, alt=7.50 — SINIRDA

# Belki gorsel cok buyuk gozukuyor — biraz kucult
for shape in slide26.shapes:
    if shape.shape_type == 13:
        # MC_Resistance gorsel — orijinal: W=6.20 H=4.80
        # Biraz kucult: H'yi 4.40 yap -> alt = 1.60+4.40 = 6.00
        shape.height = Inches(4.40)
        print('  S26: Gorsel yuksekligi kucultuldu (H=4.80 -> 4.40)')

# ====================================================
# S03 - ORTADA DURAN GORSEL DUZELTME
# Pie Class Balance — mevcut layout'a entegre et
# ====================================================
slide03 = prs.slides[2]
print('\n--- S03 layout ---')
# S03 layout (Introduction — Motivation):
# Baslik: L=0.50 T=0.45 W=12.30 H=0.55
# Subtitle: L=0.50 T=1.00 W=12.30 H=0.32
# Ana metin genellikle tam genislik (L=0.60 T=1.60 W=12.00)
# Cozum: Gorseli sag ust koseye koy, metin alaninin icine
for shape in slide03.shapes:
    L = shape.left / 914400 if shape.left else 0
    T = shape.top / 914400 if shape.top else 0
    W = shape.width / 914400 if shape.width else 0
    H = shape.height / 914400 if shape.height else 0
    if shape.shape_type == 13:
        print('  ONCEKI: L=%.2f T=%.2f W=%.2f H=%.2f' % (L, T, W, H))

# S03'e ekledigimiz gorseli cikar — bu slayt zaten Introduction, gorsel gereksiz
# Onun yerine daha uygun bir slayda ekleyecegiz
for shape in slide03.shapes:
    if shape.shape_type == 13:
        sp = shape._element
        sp.getparent().remove(sp)
        print('  S03: Gorsel cikarildi (Introduction slaytina uymuyor)')
        break

# ====================================================
# S09 - ORTADA DURAN GORSEL DUZELTME
# Architecture Flowchart — sag panele entegre et
# Layout: Sol panel (L=0.60, T=1.60, W=7.50, H=5.30) + Sag panel (L=8.30, T=1.60, W=4.60, H=5.20)
# Sag panelde: baslik (T=1.60), text (T=2.20, H=0.80), text (T=3.10, H=3.70)
# Gorsel sol panele sigmaz (metin dolu), sag panele de sigmaz (kutu dolu)
# Cozum: Gorseli cikar — bu slayt ARIMA teorik, flowchart uymuyor
# ====================================================
slide09 = prs.slides[8]
for shape in slide09.shapes:
    if shape.shape_type == 13:
        sp = shape._element
        sp.getparent().remove(sp)
        print('  S09: Gorsel cikarildi (ARIMA teorik slaytina uymuyor)')
        break

# ====================================================
# S10 - ORTADA DURAN GORSEL DUZELTME
# Search Space Specs — sag panele entegre et
# Layout: Sol panel (L=0.60, T=1.60, W=7.50) + Sag panel (L=8.30, T=1.60, W=4.60)
# Cozum: Gorseli cikar — bu slayt LSTM teorik
# Bunun yerine S14 (Search Space) slaytina ekleyecegiz
# ====================================================
slide10 = prs.slides[9]
for shape in slide10.shapes:
    if shape.shape_type == 13:
        sp = shape._element
        sp.getparent().remove(sp)
        print('  S10: Gorsel cikarildi (LSTM teorik slaytina uymuyor)')
        break

# ====================================================
# S25 - Heatmap ORTADA DURUYOR
# Layout: Tablo (T=1.60, H=1.85) + Sol gorsel (T=3.65, W=6.20) + Sag metin (T=3.65)
# Heatmap sag panelde metin uzerine binmis
# Cozum: Gorseli cikar — S25 zaten Risk_vs_Learning gorseli + tablo + metin ile dolu
# ====================================================
slide25 = prs.slides[24]
pics_25 = [s for s in slide25.shapes if s.shape_type == 13]
for shape in pics_25:
    L = shape.left / 914400 if shape.left else 0
    if L > 6.0:  # Sag taraftaki = bizim ekledigimiz heatmap
        sp = shape._element
        sp.getparent().remove(sp)
        print('  S25: Heatmap cikarildi (slayt zaten dolu)')
        break

# ====================================================
# YENIDEN EKLEME — Daha uygun slaytlara, daha iyi pozisyonlarla
# ====================================================

# S16 (MC Detection) — Sol panelde bos alan var (metin kisa ise)
# MC Sinkhole gorseli bu slayta uygun — SAG panelin altina ekle
slide16 = prs.slides[15]
# Sag panel: Rectangle L=7.60, T=1.55, W=5.30, H=5.30
# Sag panel icerik: Tablo T=2.05 H=3.00 (alt=5.05), TextBox T=5.15 H=1.60 (alt=6.75)
# Bos alan yok sag panelde. Sol panelde metin var T=1.55 H=5.30 (alt=6.85)
# Cozum: Gorseli sol panelin alt kismine KOY, metni ustune bindirmeden
# Sol panelde metin okunabilir kalsin, gorseli sol alt kosede kucuk koy
pic16 = slide16.shapes.add_picture(
    os.path.join(GRAFIK, '21_The_MC_Sinkhole.png'),
    left=Inches(0.60), top=Inches(5.00),
    width=Inches(3.50), height=Inches(2.00)
)
print('  S16: MC Sinkhole eklendi (sol alt, kucuk, W=3.50 H=2.00)')

# S15 (Evaluation Metrics) — bos slayt, gorsel eklenebilir
# Pie Class Balance bu slayta daha uygun
slide15 = prs.slides[14]
# S15 layout kontrolu
has_content_15 = False
for shape in slide15.shapes:
    if shape.shape_type == 13 or shape.has_table:
        has_content_15 = True
if not has_content_15:
    # Sag panele (genellikle L=7.5+) veya sol alt bosluga koy
    pic15 = slide15.shapes.add_picture(
        os.path.join(GRAFIK, '13_Pie_Class_Balance.png'),
        left=Inches(8.00), top=Inches(3.50),
        width=Inches(4.50), height=Inches(3.20)
    )
    print('  S15: Pie Class Balance eklendi (sag panel, W=4.50 H=3.20)')

# S14 (Search Space) — 2 tablo yan yana, ust kisimda 5 card
# Search Space Specs gorseli buraya uygun ama yer yok
# Cikaralim, 16_Search_Space zaten metin halinde tabloda var

# Sonuc olarak eklenen gorseller:
# S15: Pie Class Balance
# S16: MC Sinkhole

prs.save(DST)
print('\nKaydedildi!')

# Final dogrulama
prs2 = Presentation(DST)
total_i = 0
total_t = 0
for i, slide in enumerate(prs2.slides, 1):
    imgs = sum(1 for s in slide.shapes if s.shape_type == 13)
    tbls = sum(1 for s in slide.shapes if s.has_table)
    total_i += imgs
    total_t += tbls
    if imgs:
        for s in slide.shapes:
            if s.shape_type == 13:
                L=s.left/914400; T=s.top/914400; W=s.width/914400; H=s.height/914400
                b = T+H
                ok = 'OK' if b <= 7.20 else 'TASMA!'
                print('  S%02d: IMG L=%.2f T=%.2f W=%.2f H=%.2f alt=%.2f [%s]' % (i, L, T, W, H, b, ok))

print('\nTOPLAM: %d gorsel, %d tablo' % (total_i, total_t))
