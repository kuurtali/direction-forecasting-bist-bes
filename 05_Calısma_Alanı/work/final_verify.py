# -*- coding: utf-8 -*-
import sys, io, os
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
from pptx import Presentation

prs = Presentation(r'c:\Users\Kurt\Desktop\Proje\BES_Pension_UYIK_2026_FINAL.pptx')

print('FINAL DOGRULAMA:')
print('Toplam slayt: ' + str(len(prs.slides)))
total_i = 0
total_t = 0
for i, slide in enumerate(prs.slides, 1):
    imgs = sum(1 for s in slide.shapes if s.shape_type == 13)
    tbls = sum(1 for s in slide.shapes if s.has_table)
    total_i += imgs
    total_t += tbls
    if imgs or tbls:
        details = []
        for s in slide.shapes:
            if s.shape_type == 13:
                l = s.left/914400; t = s.top/914400; w = s.width/914400
                details.append('IMG L=%.1f T=%.1f W=%.1f' % (l, t, w))
            if s.has_table:
                l = s.left/914400; t = s.top/914400
                details.append('TBL L=%.1f T=%.1f' % (l, t))
        joined = ' | '.join(details)
        print('  S%02d: %dg %dt | %s' % (i, imgs, tbls, joined))

print('\nTOPLAM: %d gorsel, %d tablo' % (total_i, total_t))
print('Gorsel/slayt orani: %.2f' % (total_i/len(prs.slides)))
