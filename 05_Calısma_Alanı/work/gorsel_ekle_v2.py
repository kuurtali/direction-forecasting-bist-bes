# -*- coding: utf-8 -*-
"""Temiz gorsel ekleme — BACKUP'tan baslayarak"""
import sys, io, os, shutil
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
from pptx import Presentation
from pptx.util import Inches

BASE = r'c:\Users\Kurt\Desktop\Proje'
BACKUP = os.path.join(BASE, 'BES_Pension_UYIK_2026_FINAL_BACKUP.pptx')
DST = os.path.join(BASE, 'BES_Pension_UYIK_2026_FINAL.pptx')
GRAFIK = os.path.join(BASE, '04_Gorsel_Portfolyo')

# Backup'tan geri yukle
shutil.copy2(BACKUP, DST)
print("BACKUP'tan geri yuklendi")

prs = Presentation(DST)

# S09 (idx=8): ARIMA — sag panelde info kutusu var (L=8.3, T=1.6, W=4.6, H=5.2)
# Gorseli sol panelin (L=0.6, T=1.6, W=7.5) alt bosluguna koy
slide09 = prs.slides[8]
pic09 = slide09.shapes.add_picture(
    os.path.join(GRAFIK, '11_Architecture_Flowchart.png'),
    left=Inches(0.8), top=Inches(4.5), width=Inches(7.0)
)
print("S09: Architecture Flowchart eklendi (sol alt)")

# S12 (idx=11): TEFAS Data — tablo (T=1.9, H=1.5) + 6 aciklama satiri (T=4.1-6.8)
# Cok dolu, gorseli tablonun yanina kucuk koyalim — YOK, cok dolu
# Bunun yerine S03 Introduction Motivation'a ekleyelim (bos slayt)
slide03 = prs.slides[2]
pic03 = slide03.shapes.add_picture(
    os.path.join(GRAFIK, '13_Pie_Class_Balance.png'),
    left=Inches(7.5), top=Inches(2.0), width=Inches(5.0)
)
print("S03: Pie Class Balance eklendi (Introduction — Motivation)")

# S16 (idx=15): MC Detection — sol panel metin (L=0.6, W=6.8), sag panel tablo+metin (L=7.6)
# Gorseli sol panelin alt tarafina
slide16 = prs.slides[15]
pic16 = slide16.shapes.add_picture(
    os.path.join(GRAFIK, '21_The_MC_Sinkhole.png'),
    left=Inches(0.6), top=Inches(4.5), width=Inches(6.5)
)
print("S16: MC Sinkhole eklendi (sol alt)")

# S25 (idx=24): Cross-Fund — mevcut gorsel (L=0.6, T=3.6, W=6.2) + sag panel metin (L=7.0, T=3.7)
# Heatmap'i sag panelin alt bosluguna
slide25 = prs.slides[24]
pic25 = slide25.shapes.add_picture(
    os.path.join(GRAFIK, 'EN_Graphics', '03_Model_Asset_Heatmap.png'),
    left=Inches(7.2), top=Inches(4.8), width=Inches(5.5)
)
print("S25: Model Asset Heatmap eklendi (sag alt)")

# S10 (idx=9): LSTM Framework — sag panel var ama metin agirlikli
# Buraya da bir gorsel ekleyelim
slide10 = prs.slides[9]
pic10 = slide10.shapes.add_picture(
    os.path.join(GRAFIK, '16_Search_Space_Specs.png'),
    left=Inches(0.8), top=Inches(4.5), width=Inches(7.0)
)
print("S10: Search Space Specs eklendi (LSTM Framework sol alt)")

prs.save(DST)
print(f"\nKaydedildi: {DST}")

# Dogrulama
prs2 = Presentation(DST)
total = 0
for i, slide in enumerate(prs2.slides, 1):
    imgs = sum(1 for s in slide.shapes if s.shape_type == 13)
    total += imgs
    if imgs:
        print(f"  S{i:02d}: {imgs} gorsel")
print(f"TOPLAM: {total} (onceki: 11)")
