# -*- coding: utf-8 -*-
import sys, io, os
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
from pptx import Presentation

BASE = r'c:\Users\Kurt\Desktop\Proje'

files = [
    ('UYIK_1', os.path.join(BASE, 'BES_Pension_UYIK_1.pptx')),
    ('UYIK_2026_FINAL', os.path.join(BASE, 'BES_Pension_UYIK_2026_FINAL.pptx')),
    ('Referans (01/OKE)', os.path.join(BASE, '01_Savunma_ve_Ana_Metinler', 'BES_Pension_Funds_Presentation_EN_OKE_20.04.26_UYIK.pptx')),
]

d = os.path.join(BASE, '05_Hoca_Onerlileri_Ve_Ornekler')
for f in os.listdir(d):
    if f.endswith('.pptx'):
        files.append(('Ornek (05/Gizem)', os.path.join(d, f)))
        break

for name, fp in files:
    prs = Presentation(fp)
    w = prs.slide_width / 914400
    h = prs.slide_height / 914400
    
    all_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    t = p.text.strip()
                    if t: all_text.append(t.lower())
    
    s1_texts = []
    for shape in prs.slides[0].shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                t = p.text.strip()
                if t: s1_texts.append(t)
    
    s2_texts = []
    if len(prs.slides) > 1:
        for shape in prs.slides[1].shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    t = p.text.strip()
                    if t: s2_texts.append(t.lower())
    
    last_texts = []
    for shape in prs.slides[-1].shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                t = p.text.strip()
                if t: last_texts.append(t.lower())

    s1_lower = [t.lower() for t in s1_texts]
    
    checks = {
        'Yazar': any('demir' in t or 'kurt' in t or 'ozguden' in t or 'özgüden' in t for t in s1_lower),
        'Danısman': any('karadağ' in t or 'karadag' in t or 'erdemir' in t for t in s1_lower),
        'Universite': any('hacettepe' in t for t in s1_lower),
        'Kongre adı': any('uyik' in t or 'congress' in t or 'kongre' in t for t in s1_lower),
        'Outline (S2)': any('outline' in t or 'içerik' in t for t in s2_texts),
        'References': any('reference' in t or 'kaynakça' in t or 'kaynakca' in t for t in all_text),
        'Thank You': any('thank' in t or 'teşekkür' in t for t in last_texts),
        'Lit.Review': any('literature' in t or 'literatür' in t for t in all_text),
        'Methodology': any('methodology' in t or 'yöntem' in t or 'method' in t for t in all_text),
        'Results': any('result' in t or 'sonuç' in t or 'bulgu' in t for t in all_text),
        'Conclusion': any('conclusion' in t or 'future' in t or 'sonuç' in t for t in all_text),
    }
    
    print(f'=== {name} ===')
    print(f'  Boyut: {w:.2f}" x {h:.2f}"')
    print(f'  Slayt: {len(prs.slides)}')
    print(f'  Kapak: {s1_texts[0][:60] if s1_texts else "BOS"}')
    for k, v in checks.items():
        status = "OK" if v else "YOK/EKSIK"
        print(f'  {k:20s}: {status}')
    print()

# Hakem onerileri uyum kontrolu
print("=" * 70)
print("HAKEM ONERILERi UYUM KONTROLU")
print("=" * 70)
print("""
Hoca onerileri (Inceleme-Oneriler-09.04.26.docx):
  1. "Sonuclari tablolara aktaralim" -> Sunumlarda tablo var mi?
  2. "THYAO verisinde birine karar verelim" -> UYIK BES sunumunda THYAO olmamali
  3. "Emeklilik Fonu haftalik mi gunluk mu?" -> Aciklanmis mi?
  4. "Eksik gozelerde tahmin yontemi" -> Aciklanmis mi?
  5. "Literatur taramasini birlestirelim" -> Tek blok halinde mi?
  6. "APA stili kaynakca" -> APA formatinda mi?
  
RAPOR FORMATI onerileri:
  - Baslik / Ozet / Yazarlar / Anahtar Kelimeler
  - 1. GIRIS + 1.1 Literatur Taramasi
  - 2. YONTEM (2.1 ARIMA, 2.2.1 CNN, 2.2.2 LSTM, 2.3 Metrikler)
  - 3. UYGULAMA (3.1 THYAO, 3.2 BES Fonlari)
  - 4. SONUC VE BULGULAR
  - KAYNAKLAR
""")

# Kontrol edelim
for name, fp in files[:2]:  # sadece UYIK_1 ve FINAL
    prs = Presentation(fp)
    all_text_full = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    t = p.text.strip()
                    if t: all_text_full.append(t)
    
    full = ' '.join(all_text_full).lower()
    
    print(f'\n--- {name} ---')
    print(f'  THYAO referansi var mi: {"EVET - SORUN!" if "thyao" in full else "YOK - DOGRU"}')
    print(f'  APA kaynakca: {"KONTROL GEREKLI" if "reference" in full else "KAYNAK YOK"}')
    
    # ARIMA aciklanmis mi?
    print(f'  ARIMA bahsi: {"VAR" if "arima" in full else "YOK"}')
    print(f'  CNN bahsi: {"VAR" if "cnn" in full else "YOK"}')
    print(f'  LSTM bahsi: {"VAR" if "lstm" in full else "YOK"}')
    print(f'  Haftalik/Gunluk aciklama: {"VAR" if "weekly" in full or "haftalik" in full or "daily" in full or "günlük" in full else "KONTROL GEREKLI"}')
    print(f'  Missing data aciklama: {"VAR" if "missing" in full or "eksik" in full or "interpolat" in full else "KONTROL GEREKLI"}')
    print(f'  Metrik formulleri: {"VAR" if "accuracy" in full or "sensitivity" in full or "specificity" in full else "YOK"}')
