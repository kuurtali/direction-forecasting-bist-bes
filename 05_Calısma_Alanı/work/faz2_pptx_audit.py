# -*- coding: utf-8 -*-
"""Faz 2: PPTX Analizi + Faz 3: Grafik/Tablo Yeterliliği"""
from pptx import Presentation
from pptx.util import Inches
from docx import Document
import os, re, sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

BASE = r"c:\Users\Kurt\Desktop\Proje"

def analyze_pptx(filepath):
    """Analyze a PPTX file: slides, text, images, tables"""
    prs = Presentation(filepath)
    fname = os.path.basename(filepath)
    print(f"\n{'='*70}")
    print(f"PPTX ANALİZİ: {fname}")
    print(f"  Slayt sayısı: {len(prs.slides)}")
    print(f"  Boyut: {os.path.getsize(filepath):,} bytes")
    print(f"{'='*70}")
    
    total_images = 0
    total_tables = 0
    total_charts = 0
    all_numbers = []
    
    for i, slide in enumerate(prs.slides, 1):
        slide_images = 0
        slide_tables = 0
        slide_text = []
        
        for shape in slide.shapes:
            # Text
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_text.append(text)
            
            # Tables
            if shape.has_table:
                slide_tables += 1
                total_tables += 1
                tbl = shape.table
                print(f"\n  Slayt {i} - TABLO ({tbl.rows.__len__()} satır × {len(tbl.columns)} sütun):")
                for row_idx, row in enumerate(tbl.rows):
                    cells = [cell.text.strip() for cell in row.cells]
                    if row_idx < 3 or any(c for c in cells if c):  # İlk 3 satır + dolu satırlar
                        print(f"    Satır {row_idx}: {cells}")
                    if row_idx == 3 and len(tbl.rows) > 5:
                        print(f"    ... ({len(tbl.rows)-4} satır daha)")
            
            # Images
            if shape.shape_type == 13:  # Picture
                slide_images += 1
                total_images += 1
                img = shape.image
                print(f"  Slayt {i} - GÖRSEl: {img.content_type}, {len(img.blob):,} bytes")
        
        # Print slide summary
        title = ""
        for t in slide_text[:1]:
            title = t[:80]
        print(f"\n  Slayt {i}: \"{title}\"")
        if slide_images: print(f"    → {slide_images} görsel")
        if slide_tables: print(f"    → {slide_tables} tablo")
        
        # Extract numbers for verification
        for t in slide_text:
            nums = re.findall(r'\d+\.\d+%?|\d+%', t)
            if nums:
                all_numbers.extend([(i, n) for n in nums])
    
    print(f"\n{'='*70}")
    print(f"ÖZET — {fname}:")
    print(f"  Toplam slayt: {len(prs.slides)}")
    print(f"  Toplam görsel: {total_images}")
    print(f"  Toplam tablo: {total_tables}")
    print(f"{'='*70}")
    
    return {
        'slides': len(prs.slides),
        'images': total_images,
        'tables': total_tables,
        'numbers': all_numbers
    }

# === PPTX 1 ===
pptx1 = os.path.join(BASE, "BES_Pension_UYIK_1.pptx")
result1 = analyze_pptx(pptx1)

# === PPTX 2 ===
pptx2 = os.path.join(BASE, "BES_Pension_UYIK_2026_FINAL.pptx")
result2 = analyze_pptx(pptx2)

# === Referans PPTX (01 içindeki) ===
pptx_ref = os.path.join(BASE, "01_Savunma_ve_Ana_Metinler", "BES_Pension_Funds_Presentation_EN_OKE_20.04.26_UYIK.pptx")
if os.path.exists(pptx_ref):
    result_ref = analyze_pptx(pptx_ref)

# === FAZ 3B: DOC1/2/3 Tablo Analizi ===
print(f"\n{'='*70}")
print("FAZ 3B: DOC1/2/3 TABLO ANALİZİ")
print(f"{'='*70}")

doc_files = [
    "DOC1 Terimler Formuller.docx",
    "DOC2 Appendix Koleksiyonu.docx",
    "DOC3 Model Karsilastirma Tablolari.docx"
]

for docname in doc_files:
    docpath = os.path.join(BASE, "01_Savunma_ve_Ana_Metinler", docname)
    if not os.path.exists(docpath):
        print(f"\n  {docname}: DOSYA BULUNAMADI!")
        continue
    doc = Document(docpath)
    print(f"\n  === {docname} ===")
    print(f"  Paragraf sayısı: {len(doc.paragraphs)}")
    print(f"  Tablo sayısı: {len(doc.tables)}")
    
    # İlk 3 paragrafı göster (başlık bilgisi)
    for p in doc.paragraphs[:5]:
        if p.text.strip():
            print(f"    Paragraf: \"{p.text.strip()[:100]}\"")
    
    # Tablo detayları
    for t_idx, table in enumerate(doc.tables):
        rows = len(table.rows)
        cols = len(table.columns)
        print(f"\n    Tablo {t_idx+1}: {rows} satır × {cols} sütun")
        # Başlık satırı
        if rows > 0:
            header = [cell.text.strip()[:30] for cell in table.rows[0].cells]
            print(f"      Başlık: {header}")
        # İlk veri satırı
        if rows > 1:
            first_row = [cell.text.strip()[:30] for cell in table.rows[1].cells]
            print(f"      İlk satır: {first_row}")

# === FAZ 3A: 04 Grafik Karşılaştırması ===
print(f"\n{'='*70}")
print("FAZ 3A: GRAFİK YETERLİLİĞİ KARŞILAŞTIRMASI")
print(f"{'='*70}")

grafik_dir = os.path.join(BASE, "04_Gorsel_Portfolyo")
all_pngs = []
for root, dirs, files in os.walk(grafik_dir):
    for f in files:
        if f.endswith('.png'):
            rel = os.path.relpath(os.path.join(root, f), grafik_dir)
            all_pngs.append(rel)

print(f"\n  04_Gorsel_Portfolyo'daki toplam PNG: {len(all_pngs)}")
for p in sorted(all_pngs):
    print(f"    {p}")

print(f"\n  PPTX'lerdeki toplam gömülü görsel:")
print(f"    BES_Pension_UYIK_1.pptx: {result1['images']} görsel")
print(f"    BES_Pension_UYIK_2026_FINAL.pptx: {result2['images']} görsel")

# Yeterlilik değerlendirmesi
print(f"\n{'='*70}")
print("SONUÇ DEĞERLENDİRMESİ")
print(f"{'='*70}")

min_graphics = 8
min_tables = 3

for name, r in [("UYIK_1", result1), ("UYIK_2026_FINAL", result2)]:
    g_ok = "✓ YETERLİ" if r['images'] >= min_graphics else f"✗ YETERSİZ (min {min_graphics} gerekli)"
    t_ok = "✓ YETERLİ" if r['tables'] >= min_tables else f"✗ YETERSİZ (min {min_tables} gerekli)"
    print(f"\n  {name}:")
    print(f"    Grafik: {r['images']} → {g_ok}")
    print(f"    Tablo:  {r['tables']} → {t_ok}")

print(f"\n=== SCRIPT TAMAMLANDI ===")
