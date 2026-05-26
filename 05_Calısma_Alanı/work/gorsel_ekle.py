# -*- coding: utf-8 -*-
"""UYIK_2026_FINAL'a 04'ten BES-uyumlu görsel ekleme"""
import sys, io, os, shutil
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
from pptx import Presentation
from pptx.util import Inches, Emu

BASE = r'c:\Users\Kurt\Desktop\Proje'
SRC = os.path.join(BASE, 'BES_Pension_UYIK_2026_FINAL.pptx')
BACKUP = os.path.join(BASE, 'BES_Pension_UYIK_2026_FINAL_BACKUP.pptx')
GRAFIK = os.path.join(BASE, '04_Gorsel_Portfolyo')

# Yedek al
shutil.copy2(SRC, BACKUP)
print(f"Yedek alindi: {BACKUP}")

prs = Presentation(SRC)
slide_w = prs.slide_width   # 13.33" = 12192000 EMU
slide_h = prs.slide_height  # 7.50" = 6858000 EMU

print(f"Slayt boyutu: {slide_w/914400:.2f}\" x {slide_h/914400:.2f}\"")

# Mevcut durum
for i, slide in enumerate(prs.slides, 1):
    imgs = sum(1 for s in slide.shapes if s.shape_type == 13)
    tbls = sum(1 for s in slide.shapes if s.has_table)
    if imgs or tbls:
        print(f"  S{i:02d}: {imgs} gorsel, {tbls} tablo")

# Gorsel ekleme plani
additions = [
    {
        'slide_idx': 8,   # S09 - Theoretical Framework — ARIMA (0-indexed)
        'image': os.path.join(GRAFIK, '11_Architecture_Flowchart.png'),
        'desc': 'Architecture Flowchart -> S09 (ARIMA)',
        # Sag alt kosede, metin alanini engellemeyecek sekilde
        'left': Inches(7.5),
        'top': Inches(1.8),
        'width': Inches(5.2),
    },
    {
        'slide_idx': 11,  # S12 - Methodology — TEFAS Data & Fund Overview
        'image': os.path.join(GRAFIK, '13_Pie_Class_Balance.png'),
        'desc': 'Pie Class Balance -> S12 (TEFAS Data)',
        'left': Inches(8.0),
        'top': Inches(2.0),
        'width': Inches(4.8),
    },
    {
        'slide_idx': 13,  # S14 - Methodology — Model Search Space
        'image': os.path.join(GRAFIK, '16_Search_Space_Specs.png'),
        'desc': 'Search Space Specs -> S14 (Search Space)',
        'left': Inches(7.0),
        'top': Inches(2.0),
        'width': Inches(5.5),
    },
    {
        'slide_idx': 15,  # S16 - Methodology — Majority Class Detection
        'image': os.path.join(GRAFIK, '21_The_MC_Sinkhole.png'),
        'desc': 'MC Sinkhole -> S16 (MC Detection)',
        'left': Inches(6.5),
        'top': Inches(1.5),
        'width': Inches(6.3),
    },
    {
        'slide_idx': 24,  # S25 - Results — Cross-Fund Champions
        'image': os.path.join(GRAFIK, 'EN_Graphics', '03_Model_Asset_Heatmap.png'),
        'desc': 'Model Asset Heatmap -> S25 (Cross-Fund)',
        'left': Inches(7.5),
        'top': Inches(1.8),
        'width': Inches(5.2),
    },
]

# Slaytlardaki mevcut icerik alanlarini kontrol et
for add in additions:
    idx = add['slide_idx']
    slide = prs.slides[idx]
    
    # Mevcut shape'lerin konumlarini kontrol et
    existing_shapes = []
    for shape in slide.shapes:
        existing_shapes.append({
            'name': shape.name[:30],
            'left': shape.left / 914400 if shape.left else 0,
            'top': shape.top / 914400 if shape.top else 0,
            'width': shape.width / 914400 if shape.width else 0,
            'height': shape.height / 914400 if shape.height else 0,
        })
    
    print(f"\nS{idx+1} ({add['desc']}):")
    for s in existing_shapes:
        print(f"  {s['name']:30s} L={s['left']:.1f} T={s['top']:.1f} W={s['width']:.1f} H={s['height']:.1f}")

print("\n" + "="*50)

# Simdi slayt icerigine gore akilli pozisyonlama
# Her slaytdaki yer durumuna gore gorsel konumlandirma
for add in additions:
    idx = add['slide_idx']
    slide = prs.slides[idx]
    img_path = add['image']
    
    if not os.path.exists(img_path):
        print(f"HATA: {img_path} bulunamadi!")
        continue
    
    # Mevcut icerik alanini bul
    max_right = 0
    content_shapes = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            right = (shape.left + shape.width) / 914400 if shape.left and shape.width else 0
            if right > max_right:
                max_right = right
            content_shapes.append(shape)
    
    # Eger slaytda tablo varsa, gorseli tablonun yanina yerlestir
    has_table = any(s.has_table for s in slide.shapes)
    
    if has_table:
        # Tablo olan slaytlarda daha kucuk gorsel, sag tarafa
        left = Inches(8.5)
        top = Inches(2.0)
        width = Inches(4.3)
    else:
        # Tablo olmayan slaytlarda daha buyuk gorsel, ortaya yakin
        left = add['left']
        top = add['top']
        width = add['width']
    
    # Gorseli ekle
    pic = slide.shapes.add_picture(img_path, left, top, width=width)
    print(f"EKLENDI: {add['desc']}")
    print(f"  Konum: L={left/914400:.1f}\" T={top/914400:.1f}\" W={width/914400:.1f}\"")

# Kaydet
prs.save(SRC)
print(f"\nKaydedildi: {SRC}")

# Dogrulama
prs2 = Presentation(SRC)
total_imgs = 0
for i, slide in enumerate(prs2.slides, 1):
    imgs = sum(1 for s in slide.shapes if s.shape_type == 13)
    total_imgs += imgs
    if imgs:
        print(f"  S{i:02d}: {imgs} gorsel")
print(f"\nTOPLAM GORSEL: {total_imgs} (onceki: 11)")
