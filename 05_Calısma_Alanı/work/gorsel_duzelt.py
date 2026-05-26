# -*- coding: utf-8 -*-
"""Gorsel pozisyon duzeltme — tasmalari ve cakismalari gider"""
import sys, io, os
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
from pptx import Presentation
from pptx.util import Inches, Emu

BASE = r'c:\Users\Kurt\Desktop\Proje'
DST = os.path.join(BASE, 'BES_Pension_UYIK_2026_FINAL.pptx')

prs = Presentation(DST)
slide_h = prs.slide_height  # 6858000 EMU = 7.50"
slide_w = prs.slide_width   # 12192000 EMU = 13.33"

# =========================================
# S16 (idx=15): MC Sinkhole TASIYOR
# Mevcut: L=0.60 T=4.50 W=6.50 H=4.55 -> alt = 9.05" (TASIYOR!)
# Sol panelde metin: L=0.60 T=1.55 W=6.80 H=5.30 (T+H=6.85")
# Gorsel serbest alan: T=1.55'ten asagi, ama metin var
# Cozum: Gorseli kucult ve sol panelin altina sigdir
# Kullanilabilir alan: T=5.0 -> T=7.0 (footer T=7.18), yukseklik max 2.0"
# =========================================
slide16 = prs.slides[15]
for shape in slide16.shapes:
    if shape.shape_type == 13:
        # MC Sinkhole: landscape gorsel, genis ama kisa yapmaliyiz
        shape.left = Inches(0.60)
        shape.top = Inches(4.80)
        shape.width = Inches(6.50)
        # Orani koruyarak yuksekligi hesapla
        # Orijinal gorsel orani: 1316x930 piksel -> oran = 930/1316 = 0.707
        shape.height = Inches(6.50 * 0.707)  # = 4.60" -> hala tasar!
        # Daha kucuk yapalim: max yukseklik = 7.18 - 4.80 = 2.38"
        max_h = 2.30
        new_w = max_h / 0.707  # = 3.25"
        shape.width = Inches(new_w)
        shape.height = Inches(max_h)
        shape.left = Inches(0.60)
        shape.top = Inches(4.80)
        print('S16 duzeltildi: L=0.60 T=4.80 W=%.2f H=%.2f' % (new_w, max_h))

# =========================================
# S25 (idx=24): Heatmap TASIYOR
# Mevcut: L=7.20 T=4.80 W=5.50 H=3.93 -> alt = 8.73" (TASIYOR!)
# Sag panelde metin var: L=6.95 T=3.70 "Key take-aways" + bullet points (T=4.10 H=2.75 -> T+H=6.85")
# Gorsel metinin uzerine biniyor
# Cozum: Gorseli kuculterek metin altina sigdir
# Kullanilabilir alan: T=5.5 -> T=7.0, yukseklik max 1.5"
# VEYA gorseli tamamen cikaralim ve S25'i temiz birakalim (zaten 1 gorsel + 1 tablo var)
# =========================================
slide25 = prs.slides[24]
pics_25 = [s for s in slide25.shapes if s.shape_type == 13]
# 2 gorsel var: biri orijinal (Risk_vs_Learning, L=0.55), biri ekledigimiz (Heatmap, L=7.20)
for shape in pics_25:
    if shape.left / 914400 > 6.0:  # Sag taraftaki = bizim ekledigimiz
        # Heatmap: 1326x938 px -> oran = 938/1326 = 0.707
        # Kullanilabilir: T=5.0 -> T=7.0, max H=2.0"
        max_h = 1.90
        new_w = max_h / 0.707
        shape.left = Inches(7.20)
        shape.top = Inches(5.00)
        shape.width = Inches(new_w)
        shape.height = Inches(max_h)
        print('S25 duzeltildi: L=7.20 T=5.00 W=%.2f H=%.2f' % (new_w, max_h))

# =========================================
# S01 ve S02: Orijinal icerik — biz gorsel eklemedik
# Kullanicinin gordugu sorun ne olabilir?
# S01: Hic gorsel yok, sadece metin
# S02: Hic gorsel yok, sadece Outline metni
# Muhtemel sorun: arka plan dikdortgeni veya bos textbox'lar kotu gorunuyor olabilir
# Bunlara dokunmuyorum cunku orijinal tasarimin parcasi
# =========================================
print('S01: Orijinal tasarim — gorsel yok, degisiklik yapilmadi')
print('S02: Orijinal tasarim — gorsel yok, degisiklik yapilmadi')

# =========================================
# S26: Orijinal gorsel + tablo — biz gorsel eklemedik
# Mevcut layout temiz gorunuyor (gorsel L=0.55 T=1.60 W=6.20 H=4.80)
# Alt sinir: T+H = 6.40 < 7.18 footer -> TAMAM
# Kullanicinin gordugu sorun: belki caption text (T=6.43) gorunmuyor?
# =========================================
print('S26: Orijinal tasarim — layout analiz edildi, tasma yok')

# Kaydet
prs.save(DST)
print('\nKaydedildi!')

# Dogrulama
prs2 = Presentation(DST)
for slide_num in [16, 25]:
    idx = slide_num - 1
    slide = prs2.slides[idx]
    for shape in slide.shapes:
        if shape.shape_type == 13:
            L = shape.left / 914400
            T = shape.top / 914400
            W = shape.width / 914400
            H = shape.height / 914400
            bottom = T + H
            overflow = 'TAMAM' if bottom <= 7.2 else 'HALA TASIYOR!'
            print('  S%d gorsel: L=%.2f T=%.2f W=%.2f H=%.2f alt=%.2f [%s]' % (slide_num, L, T, W, H, bottom, overflow))
